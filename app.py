import streamlit as st
import json
import pypdf
import pdfplumber
from pptx import Presentation
import os
from groq import Groq
from dotenv import load_dotenv
import tempfile
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch

# Load environment variables
load_dotenv()

# Initialize Groq client
def get_groq_client():
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        st.error("Please set your GROQ_API_KEY in the environment variables or .env file")
        st.stop()
    return Groq(api_key=api_key)

def extract_text_from_pdf(file_path):
    """Enhanced PDF text extraction using multiple methods for better coverage"""
    text = ""
    
    # Method 1: Try pdfplumber first (better for complex layouts)
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                
                # Also extract text from tables if any
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        if row:
                            text += " | ".join([cell or "" for cell in row]) + "\n"
        
        if text.strip():
            return text
    except Exception as e:
        st.warning(f"pdfplumber extraction failed: {str(e)}, trying pypdf...")
    
    # Method 2: Fallback to pypdf if pdfplumber fails
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = pypdf.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        
        if text.strip():
            return text
    except Exception as e:
        st.error(f"Error extracting text from PDF with both methods: {str(e)}")
        return None
    
    return text if text.strip() else None

def extract_text_from_pptx(file_path):
    """Enhanced PPTX text extraction including slide notes and comprehensive content"""
    try:
        prs = Presentation(file_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text += " | ".join(row_text) + "\n"
                
                # Extract text from text frames and paragraphs
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text += paragraph.text + "\n"
            
            # Extract slide notes if available
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame'):
                    notes_text = notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        text += f"[Slide Notes: {notes_text}]\n"
        
        return text.strip() if text.strip() else None
        
    except Exception as e:
        st.error(f"Error extracting text from PPTX: {str(e)}")
        return None

def evaluate_submission(extracted_text):
    """Enhanced evaluation with advanced features using Groq API"""
    client = get_groq_client()
    
    prompt = f"""You are an expert AI evaluator for pitch decks, MSME proposals, and hackathon submissions. Analyze the extracted content and provide a comprehensive evaluation with advanced features.

🎯 Evaluation Criteria (Each scored from 1 to 10):
1. **Clarity** – Is the idea, problem, and solution clearly explained?
2. **Innovation** – How novel, original, or disruptive is the idea?
3. **Feasibility** – Is the solution technically and practically implementable?
4. **Presentation Quality** – Is the content professional, structured, and engaging?
5. **Impact** – What is the expected market, economic, or social impact?
6. **Theme Alignment** – How well does it align with competition goals or MSME objectives?

🔍 Advanced Analysis Required:
1. **Detected Theme/Domain** (e.g., FinTech, HealthTech, Sustainability, EdTech, AgriTech)
2. **Top Keywords** (5-10 core concepts or buzzwords from the content)
3. **Suggested Project Title** (Max 10 words, catchy and relevant)
4. **2-Line Summary** (Concise description of the idea)
5. **Improvement Suggestions** (3 concrete actionable points)
6. **Recommended Resources** (2 relevant tools, frameworks, or platforms)
7. **Visual Quality Check** (Assess design/structure quality from textual cues)
8. **Pitch Readiness Score** (Out of 10 – readiness for investors/juries)

📄 Submission Content:
\"\"\"{extracted_text}\"\"\"

✅ Output Format (Return ONLY valid JSON):
{{
    "scores": {{
        "clarity": 0,
        "innovation": 0,
        "feasibility": 0,
        "presentation": 0,
        "impact": 0,
        "theme_alignment": 0
    }},
    "total_score": 0,
    "grade": "A+",
    "feedback_summary": "Professional 3-line summary of the submission's strengths and areas for improvement.",
    "theme": "Primary domain/industry category",
    "keywords": ["keyword1", "keyword2", "keyword3", "keyword4", "keyword5"],
    "project_title": "Catchy Project Title (Max 10 Words)",
    "project_summary": "Two-line concise description of the core idea and its value proposition.",
    "improvement_suggestions": [
        "Specific actionable improvement 1",
        "Specific actionable improvement 2", 
        "Specific actionable improvement 3"
    ],
    "recommended_resources": [
        "Relevant tool/framework/platform 1",
        "Relevant tool/framework/platform 2"
    ],
    "visual_quality_comment": "Assessment of presentation design and structure quality based on content organization",
    "pitch_readiness_score": 0
}}"""

    try:
        response = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
            max_tokens=1500  # Increased token limit for comprehensive response
        )
        
        result_text = response.choices[0].message.content.strip()
        
        # Debug: Show raw response in expander for troubleshooting
        with st.expander("🔧 Debug: Raw AI Response", expanded=False):
            st.text_area("Raw Response", result_text, height=200)
        
        # Enhanced JSON extraction with multiple fallback methods
        json_text = result_text
        
        # Method 1: Extract from ```json blocks
        if "```json" in result_text:
            json_text = result_text.split("```json")[1].split("```")[0].strip()
        # Method 2: Extract from ``` blocks
        elif "```" in result_text:
            json_text = result_text.split("```")[1].strip()
        # Method 3: Look for JSON-like structure
        elif "{" in result_text and "}" in result_text:
            start = result_text.find("{")
            end = result_text.rfind("}") + 1
            json_text = result_text[start:end]
        
        # Clean up common issues
        json_text = json_text.strip()
        
        # Try to parse JSON
        if json_text:
            try:
                parsed_result = json.loads(json_text)
                
                # Validate required fields and provide defaults if missing
                default_result = {
                    "scores": {
                        "clarity": parsed_result.get("scores", {}).get("clarity", 5),
                        "innovation": parsed_result.get("scores", {}).get("innovation", 5),
                        "feasibility": parsed_result.get("scores", {}).get("feasibility", 5),
                        "presentation": parsed_result.get("scores", {}).get("presentation", 5),
                        "impact": parsed_result.get("scores", {}).get("impact", 5),
                        "theme_alignment": parsed_result.get("scores", {}).get("theme_alignment", 5)
                    },
                    "total_score": parsed_result.get("total_score", 30),
                    "grade": parsed_result.get("grade", "B"),
                    "feedback_summary": parsed_result.get("feedback_summary", "Evaluation completed successfully."),
                    "theme": parsed_result.get("theme", "General"),
                    "keywords": parsed_result.get("keywords", ["innovation", "technology", "solution"]),
                    "project_title": parsed_result.get("project_title", "Innovative Solution"),
                    "project_summary": parsed_result.get("project_summary", "A comprehensive solution addressing key challenges."),
                    "improvement_suggestions": parsed_result.get("improvement_suggestions", [
                        "Enhance clarity in presentation",
                        "Provide more detailed implementation plan",
                        "Include market analysis and validation"
                    ]),
                    "recommended_resources": parsed_result.get("recommended_resources", [
                        "Business Model Canvas",
                        "Lean Startup Methodology"
                    ]),
                    "visual_quality_comment": parsed_result.get("visual_quality_comment", "Content appears well-structured."),
                    "pitch_readiness_score": parsed_result.get("pitch_readiness_score", 6)
                }
                
                return default_result
                
            except json.JSONDecodeError as json_error:
                st.error(f"JSON parsing failed: {str(json_error)}")
                st.error(f"Attempted to parse: {json_text[:200]}...")
                return None
        else:
            st.error("No valid JSON content found in AI response")
            return None
    
    except Exception as e:
        st.error(f"Error calling Groq API: {str(e)}")
        return None

def generate_pdf_report(evaluation_result, filename):
    """Generate comprehensive PDF report with advanced analysis"""
    buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
    doc = SimpleDocTemplate(buffer.name, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=18,
        spaceAfter=30,
        alignment=1  # Center alignment
    )
    story.append(Paragraph("📝 Advanced Submission Evaluation Report", title_style))
    story.append(Spacer(1, 20))
    
    # File info and basic metrics
    story.append(Paragraph(f"<b>File:</b> {filename}", styles['Normal']))
    story.append(Paragraph(f"<b>Theme/Domain:</b> {evaluation_result.get('theme', 'Not detected')}", styles['Normal']))
    story.append(Paragraph(f"<b>Pitch Readiness Score:</b> {evaluation_result.get('pitch_readiness_score', 0)}/10", styles['Normal']))
    story.append(Spacer(1, 15))
    
    # Project insights section
    story.append(Paragraph("<b>💡 Project Insights</b>", styles['Heading2']))
    story.append(Paragraph(f"<b>Suggested Title:</b> {evaluation_result.get('project_title', 'Not generated')}", styles['Normal']))
    story.append(Spacer(1, 8))
    story.append(Paragraph(f"<b>Project Summary:</b>", styles['Normal']))
    story.append(Paragraph(evaluation_result.get('project_summary', 'Not available'), styles['Normal']))
    story.append(Spacer(1, 8))
    
    # Keywords
    keywords = evaluation_result.get('keywords', [])
    if keywords:
        story.append(Paragraph(f"<b>Key Concepts:</b> {', '.join(keywords)}", styles['Normal']))
    story.append(Spacer(1, 15))
    
    # Scores table
    story.append(Paragraph("<b>📊 Detailed Evaluation Scores</b>", styles['Heading2']))
    score_data = [
        ['Criteria', 'Score (out of 10)'],
        ['Clarity', str(evaluation_result['scores']['clarity'])],
        ['Innovation', str(evaluation_result['scores']['innovation'])],
        ['Feasibility', str(evaluation_result['scores']['feasibility'])],
        ['Presentation Quality', str(evaluation_result['scores']['presentation'])],
        ['Impact', str(evaluation_result['scores']['impact'])],
        ['Theme Alignment', str(evaluation_result['scores']['theme_alignment'])],
        ['', ''],
        ['Total Score', f"{evaluation_result['total_score']}/60"],
        ['Grade', evaluation_result['grade']]
    ]
    
    score_table = Table(score_data, colWidths=[3*inch, 1.5*inch])
    score_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -3), colors.beige),
        ('BACKGROUND', (0, -2), (-1, -1), colors.lightblue),
        ('FONTNAME', (0, -2), (-1, -1), 'Helvetica-Bold'),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    
    story.append(score_table)
    story.append(Spacer(1, 20))
    
    # Visual quality assessment
    story.append(Paragraph("<b>🎨 Visual Quality Assessment:</b>", styles['Heading2']))
    visual_comment = evaluation_result.get('visual_quality_comment', 'Not assessed')
    story.append(Paragraph(visual_comment, styles['Normal']))
    story.append(Spacer(1, 15))
    
    # Feedback summary
    story.append(Paragraph("<b>💬 Professional Feedback:</b>", styles['Heading2']))
    story.append(Paragraph(evaluation_result['feedback_summary'], styles['Normal']))
    story.append(Spacer(1, 15))
    
    # Improvement suggestions
    story.append(Paragraph("<b>🛠 Improvement Suggestions:</b>", styles['Heading2']))
    for i, suggestion in enumerate(evaluation_result['improvement_suggestions'], 1):
        story.append(Paragraph(f"{i}. {suggestion}", styles['Normal']))
        story.append(Spacer(1, 5))
    story.append(Spacer(1, 10))
    
    # Recommended resources
    story.append(Paragraph("<b>📚 Recommended Resources:</b>", styles['Heading2']))
    resources = evaluation_result.get('recommended_resources', [])
    for i, resource in enumerate(resources, 1):
        story.append(Paragraph(f"{i}. {resource}", styles['Normal']))
        story.append(Spacer(1, 5))
    
    doc.build(story)
    
    with open(buffer.name, 'rb') as f:
        pdf_data = f.read()
    
    os.unlink(buffer.name)
    return pdf_data

def main():
    st.set_page_config(
        page_title="📝 Submission Evaluator",
        page_icon="🎯",
        layout="wide"
    )
    
    st.title("📝 Submission Evaluator: MSME/Hackathon AI Judge")
    st.markdown("Upload your PDF or PowerPoint submission for professional evaluation")
    
    # Sidebar for API key input
    with st.sidebar:
        st.header("⚙️ Configuration")
        groq_api_key = st.text_input(
            "Groq API Key", 
            type="password", 
            help="Enter your Groq API key. Get one at https://console.groq.com/"
        )
        if groq_api_key:
            os.environ["GROQ_API_KEY"] = groq_api_key
    
    # File upload section
    st.header("📤 Upload Your Submission")
    uploaded_file = st.file_uploader(
        "Choose a PDF or PowerPoint file",
        type=['pdf', 'pptx'],
        help="Upload your pitch deck, proposal, or submission document"
    )
    
    if uploaded_file is not None:
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}") as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name
        
        try:
            # Extract text based on file type
            with st.spinner("🔍 Extracting content from your submission..."):
                if uploaded_file.name.lower().endswith('.pdf'):
                    extracted_text = extract_text_from_pdf(tmp_file_path)
                elif uploaded_file.name.lower().endswith('.pptx'):
                    extracted_text = extract_text_from_pptx(tmp_file_path)
                else:
                    st.error("Unsupported file type")
                    return
            
            if extracted_text:
                # Show extraction statistics
                word_count = len(extracted_text.split())
                char_count = len(extracted_text)
                line_count = len(extracted_text.split('\n'))
                
                st.success("✅ Content extracted successfully!")
                
                # Display extraction stats
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("📝 Words", f"{word_count:,}")
                with col2:
                    st.metric("🔤 Characters", f"{char_count:,}")
                with col3:
                    st.metric("📄 Lines", f"{line_count:,}")
                
                # Show extracted text preview with better formatting
                with st.expander("📄 Preview Extracted Content", expanded=False):
                    if len(extracted_text) > 2000:
                        st.info(f"Showing first 2000 characters of {len(extracted_text):,} total characters")
                        preview_text = extracted_text[:2000] + "\n\n... [Content truncated for preview] ..."
                    else:
                        preview_text = extracted_text
                    
                    st.text_area(
                        "Extracted Text", 
                        preview_text, 
                        height=300,
                        help="This is the text content that will be sent to AI for evaluation"
                    )
                
                # Evaluate submission
                if st.button("🎯 Evaluate Submission", type="primary"):
                    with st.spinner("🧠 AI is evaluating your submission..."):
                        evaluation_result = evaluate_submission(extracted_text)
                    
                    if evaluation_result:
                        st.success("✅ Evaluation completed!")
                        
                        # Advanced Analysis Section
                        st.header("🔍 Advanced Analysis")
                        
                        # Top row - Key insights
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("🎯 Theme/Domain", evaluation_result.get('theme', 'Not detected'))
                        with col2:
                            st.metric("🎯 Total Score", f"{evaluation_result['total_score']}/60")
                        with col3:
                            st.metric("🚀 Pitch Readiness", f"{evaluation_result.get('pitch_readiness_score', 0)}/10")
                        
                        # Project insights
                        st.subheader("💡 Project Insights")
                        col1, col2 = st.columns([1, 1])
                        
                        with col1:
                            st.write("**🏷️ Suggested Project Title:**")
                            st.info(evaluation_result.get('project_title', 'Title not generated'))
                            
                            st.write("**📝 Project Summary:**")
                            st.write(evaluation_result.get('project_summary', 'Summary not available'))
                        
                        with col2:
                            st.write("**🔑 Top Keywords:**")
                            keywords = evaluation_result.get('keywords', [])
                            if keywords:
                                # Display keywords as tags
                                keyword_html = " ".join([f"<span style='background-color: #e1f5fe; padding: 2px 8px; border-radius: 12px; margin: 2px; display: inline-block;'>{kw}</span>" for kw in keywords])
                                st.markdown(keyword_html, unsafe_allow_html=True)
                            else:
                                st.write("No keywords detected")
                        
                        # Scores section
                        st.header("📊 Detailed Evaluation Scores")
                        col1, col2 = st.columns([2, 1])
                        
                        with col1:
                            # Scores table
                            scores_data = {
                                "Criteria": ["Clarity", "Innovation", "Feasibility", "Presentation Quality", "Impact", "Theme Alignment"],
                                "Score (out of 10)": [
                                    evaluation_result['scores']['clarity'],
                                    evaluation_result['scores']['innovation'],
                                    evaluation_result['scores']['feasibility'],
                                    evaluation_result['scores']['presentation'],
                                    evaluation_result['scores']['impact'],
                                    evaluation_result['scores']['theme_alignment']
                                ]
                            }
                            
                            st.dataframe(scores_data, use_container_width=True)
                        
                        with col2:
                            # Grade with color coding
                            grade = evaluation_result['grade']
                            if grade in ['A+', 'A']:
                                st.success(f"🏅 Grade: {grade}")
                            elif grade == 'B':
                                st.warning(f"🏅 Grade: {grade}")
                            else:
                                st.error(f"🏅 Grade: {grade}")
                            
                            # Visual quality assessment
                            st.write("**🎨 Visual Quality:**")
                            visual_comment = evaluation_result.get('visual_quality_comment', 'Not assessed')
                            st.write(visual_comment)
                        
                        # Feedback and suggestions
                        st.header("💬 Professional Feedback")
                        st.info(evaluation_result['feedback_summary'])
                        
                        # Improvement suggestions and resources
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            st.subheader("🛠 Improvement Suggestions")
                            for i, suggestion in enumerate(evaluation_result['improvement_suggestions'], 1):
                                st.write(f"{i}. {suggestion}")
                        
                        with col2:
                            st.subheader("📚 Recommended Resources")
                            resources = evaluation_result.get('recommended_resources', [])
                            for i, resource in enumerate(resources, 1):
                                st.write(f"{i}. {resource}")
                        
                        # Generate PDF report
                        st.header("📥 Export Report")
                        if st.button("📄 Generate PDF Report"):
                            with st.spinner("📄 Generating PDF report..."):
                                pdf_data = generate_pdf_report(evaluation_result, uploaded_file.name)
                                st.download_button(
                                    label="⬇️ Download PDF Report",
                                    data=pdf_data,
                                    file_name=f"evaluation_report_{uploaded_file.name.split('.')[0]}.pdf",
                                    mime="application/pdf"
                                )
        
        finally:
            # Clean up temporary file
            if os.path.exists(tmp_file_path):
                os.unlink(tmp_file_path)

if __name__ == "__main__":
    main()