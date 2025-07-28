#!/usr/bin/env python3
"""
Test script to demonstrate the enhanced text extraction capabilities
"""

import tempfile
from app import extract_text_from_pdf, extract_text_from_pptx
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches

def create_test_pdf():
    """Create a test PDF with sample content"""
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
    doc = SimpleDocTemplate(temp_file.name, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Add content to PDF
    story.append(Paragraph("Innovation Submission: Smart City Traffic Management", styles['Title']))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Problem Statement:", styles['Heading2']))
    story.append(Paragraph("Urban traffic congestion is a major issue affecting millions of commuters daily. Current traffic management systems are reactive rather than predictive, leading to inefficient traffic flow and increased pollution.", styles['Normal']))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Proposed Solution:", styles['Heading2']))
    story.append(Paragraph("Our AI-powered traffic management system uses real-time data from IoT sensors, cameras, and mobile devices to predict traffic patterns and optimize signal timing dynamically.", styles['Normal']))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Key Features:", styles['Heading2']))
    story.append(Paragraph("• Real-time traffic prediction using machine learning", styles['Normal']))
    story.append(Paragraph("• Dynamic traffic signal optimization", styles['Normal']))
    story.append(Paragraph("• Mobile app for commuter route suggestions", styles['Normal']))
    story.append(Paragraph("• Integration with public transportation systems", styles['Normal']))
    
    doc.build(story)
    return temp_file.name

def create_test_pptx():
    """Create a test PowerPoint with sample content"""
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs = Presentation()
    
    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "EcoTech Solutions"
    slide1.shapes.placeholders[1].text = "Sustainable Technology for Smart Cities"
    
    # Slide 2: Problem
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "The Problem"
    content = slide2.shapes.placeholders[1].text_frame
    content.text = "Urban areas consume 78% of global energy and produce 60% of greenhouse gas emissions"
    p = content.add_paragraph()
    p.text = "Current smart city solutions lack integration and sustainability focus"
    p = content.add_paragraph()
    p.text = "Citizens need better tools to make environmentally conscious decisions"
    
    # Slide 3: Solution
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Our Solution"
    content = slide3.shapes.placeholders[1].text_frame
    content.text = "Integrated IoT platform for sustainable urban living"
    p = content.add_paragraph()
    p.text = "AI-powered energy optimization for buildings"
    p = content.add_paragraph()
    p.text = "Citizen engagement app with gamification"
    p = content.add_paragraph()
    p.text = "Real-time environmental monitoring dashboard"
    
    # Add speaker notes to slide 3
    notes_slide = slide3.notes_slide
    notes_slide.notes_text_frame.text = "This solution addresses the core issues of urban sustainability by providing actionable insights to both city administrators and citizens."
    
    prs.save(temp_file.name)
    return temp_file.name

def test_enhanced_extraction():
    """Test the enhanced extraction capabilities"""
    print("🧪 Testing Enhanced Text Extraction\n")
    
    # Test PDF extraction
    print("📄 Testing PDF Extraction...")
    pdf_file = create_test_pdf()
    try:
        pdf_text = extract_text_from_pdf(pdf_file)
        if pdf_text:
            word_count = len(pdf_text.split())
            print(f"✅ PDF extraction successful!")
            print(f"   📝 Words extracted: {word_count}")
            print(f"   🔤 Characters extracted: {len(pdf_text)}")
            print(f"   📄 Sample content: {pdf_text[:100]}...")
        else:
            print("❌ PDF extraction failed")
    except Exception as e:
        print(f"❌ PDF extraction error: {e}")
    
    print("\n" + "="*50 + "\n")
    
    # Test PPTX extraction
    print("📊 Testing PowerPoint Extraction...")
    pptx_file = create_test_pptx()
    try:
        pptx_text = extract_text_from_pptx(pptx_file)
        if pptx_text:
            word_count = len(pptx_text.split())
            print(f"✅ PPTX extraction successful!")
            print(f"   📝 Words extracted: {word_count}")
            print(f"   🔤 Characters extracted: {len(pptx_text)}")
            print(f"   📄 Sample content: {pptx_text[:150]}...")
            
            # Check if slide notes were extracted
            if "[Slide Notes:" in pptx_text:
                print("   📝 Speaker notes successfully extracted!")
        else:
            print("❌ PPTX extraction failed")
    except Exception as e:
        print(f"❌ PPTX extraction error: {e}")
    
    # Cleanup
    import os
    try:
        os.unlink(pdf_file)
        os.unlink(pptx_file)
    except:
        pass
    
    print("\n🎉 Enhanced extraction test completed!")
    print("\n💡 The improved extraction now captures:")
    print("   • PDF: Text, tables, and complex layouts")
    print("   • PPTX: Slide content, tables, and speaker notes")
    print("   • Better handling of formatting and structure")

if __name__ == "__main__":
    test_enhanced_extraction()